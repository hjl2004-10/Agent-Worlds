"""
文档处理 - 原子层 (L2)
各格式文件的文本提取，纯函数无副作用
"""
from pathlib import Path
from typing import Tuple


# 支持的文档扩展名
DOCUMENT_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.pptx', '.ppt',
    '.xlsx', '.xls', '.csv',
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp',
}


def is_document(path: Path) -> bool:
    """判断是否为需要特殊解析的文档"""
    return path.suffix.lower() in DOCUMENT_EXTENSIONS


def extract_text(path: Path, max_chars: int = 100000) -> Tuple[str, str]:
    """
    从文档中提取文本

    Returns:
        (text, format_hint): 提取的文本, 格式提示 (如 "PDF 共5页")
    """
    suffix = path.suffix.lower()

    try:
        if suffix == '.pdf':
            return _extract_pdf(path, max_chars)
        elif suffix in ('.docx', '.doc'):
            return _extract_docx(path, max_chars)
        elif suffix in ('.pptx', '.ppt'):
            return _extract_pptx(path, max_chars)
        elif suffix in ('.xlsx', '.xls'):
            return _extract_excel(path, max_chars)
        elif suffix == '.csv':
            return _extract_csv(path, max_chars)
        elif suffix in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'):
            return _extract_image_info(path)
        else:
            return "", f"不支持的格式: {suffix}"
    except ImportError as e:
        return "", f"缺少依赖库: {e}"
    except Exception as e:
        return "", f"解析失败: {e}"


def _extract_pdf(path: Path, max_chars: int) -> Tuple[str, str]:
    """提取 PDF 文本"""
    import pdfplumber

    texts = []
    total_chars = 0

    with pdfplumber.open(str(path)) as pdf:
        total_pages = len(pdf.pages)

        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text() or ""

            # 提取表格
            tables = page.extract_tables()
            table_text = ""
            for table in tables:
                for row in table:
                    cells = [str(c or "") for c in row]
                    table_text += " | ".join(cells) + "\n"

            combined = page_text
            if table_text:
                combined += "\n[表格]\n" + table_text

            if total_chars + len(combined) > max_chars:
                remaining = max_chars - total_chars
                if remaining > 0:
                    texts.append(f"--- 第 {i + 1} 页 ---\n" + combined[:remaining])
                texts.append(f"\n...(已截断，共 {total_pages} 页)")
                break

            texts.append(f"--- 第 {i + 1} 页 ---\n" + combined)
            total_chars += len(combined)

    return "\n\n".join(texts), f"PDF 共 {total_pages} 页"


def _extract_docx(path: Path, max_chars: int) -> Tuple[str, str]:
    """提取 DOCX 文本"""
    from docx import Document

    doc = Document(str(path))
    texts = []
    total_chars = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 标记标题
        if para.style and para.style.name.startswith('Heading'):
            level = para.style.name.replace('Heading ', '').replace('Heading', '1')
            try:
                level_num = int(level)
            except ValueError:
                level_num = 1
            text = "#" * level_num + " " + text

        if total_chars + len(text) > max_chars:
            texts.append(text[:max_chars - total_chars])
            texts.append("\n...(已截断)")
            break

        texts.append(text)
        total_chars += len(text)

    # 提取表格
    for table in doc.tables:
        table_lines = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_lines.append(" | ".join(cells))

        if table_lines:
            table_text = "\n[表格]\n" + "\n".join(table_lines)
            if total_chars + len(table_text) <= max_chars:
                texts.append(table_text)
                total_chars += len(table_text)

    para_count = len(doc.paragraphs)
    table_count = len(doc.tables)
    hint = f"DOCX {para_count} 段落"
    if table_count:
        hint += f", {table_count} 个表格"

    return "\n".join(texts), hint


def _extract_pptx(path: Path, max_chars: int) -> Tuple[str, str]:
    """提取 PPTX 文本"""
    from pptx import Presentation

    prs = Presentation(str(path))
    texts = []
    total_chars = 0

    for i, slide in enumerate(prs.slides):
        slide_texts = [f"--- 幻灯片 {i + 1} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))

        slide_text = "\n".join(slide_texts)

        if total_chars + len(slide_text) > max_chars:
            texts.append(slide_text[:max_chars - total_chars])
            texts.append(f"\n...(已截断，共 {len(prs.slides)} 页)")
            break

        texts.append(slide_text)
        total_chars += len(slide_text)

    return "\n\n".join(texts), f"PPTX 共 {len(prs.slides)} 页"


def _extract_excel(path: Path, max_chars: int) -> Tuple[str, str]:
    """提取 Excel 文本"""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    texts = []
    total_chars = 0
    sheet_count = len(wb.sheetnames)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_texts = [f"--- 工作表: {sheet_name} ---"]

        row_count = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(c if c is not None else "") for c in row]
            line = " | ".join(cells)
            sheet_texts.append(line)
            row_count += 1

            if total_chars + len(line) > max_chars:
                sheet_texts.append(f"...(已截断，共 {row_count}+ 行)")
                break

        sheet_text = "\n".join(sheet_texts)
        texts.append(sheet_text)
        total_chars += len(sheet_text)

        if total_chars >= max_chars:
            break

    wb.close()
    return "\n\n".join(texts), f"Excel {sheet_count} 个工作表"


def _extract_csv(path: Path, max_chars: int) -> Tuple[str, str]:
    """提取 CSV 文本"""
    import csv

    texts = []
    total_chars = 0
    row_count = 0

    with open(str(path), 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            line = " | ".join(row)
            if total_chars + len(line) > max_chars:
                texts.append(f"...(已截断，共 {row_count}+ 行)")
                break
            texts.append(line)
            total_chars += len(line)
            row_count += 1

    return "\n".join(texts), f"CSV {row_count} 行"


def _extract_image_info(path: Path) -> Tuple[str, str]:
    """提取图片基本信息 (不做 OCR，返回元数据)"""
    from PIL import Image

    img = Image.open(str(path))
    width, height = img.size
    mode = img.mode
    fmt = img.format or path.suffix.upper().lstrip('.')

    info_lines = [
        f"[图片文件]",
        f"格式: {fmt}",
        f"尺寸: {width} x {height}",
        f"模式: {mode}",
        f"文件大小: {path.stat().st_size / 1024:.1f} KB",
    ]

    # EXIF 信息
    exif = img.getexif()
    if exif:
        interesting_tags = {
            271: "设备制造商",
            272: "设备型号",
            306: "拍摄时间",
            274: "方向",
        }
        for tag_id, label in interesting_tags.items():
            if tag_id in exif:
                info_lines.append(f"{label}: {exif[tag_id]}")

    img.close()

    return "\n".join(info_lines), f"图片 {width}x{height} {fmt}"
